from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.dml.color import RGBColor
import uuid
import re
import os
import json
import toml
from datetime import datetime

# 加载配置文件
config_path = os.path.join(os.path.dirname(__file__), 'config.toml')
with open(config_path, 'r') as f:
    config = toml.load(f)

workspace_path = config['workspace']['path']

# 如果不存在，则创建目录
if not os.path.exists(workspace_path):
    os.makedirs(workspace_path)

def markdown_to_hierarchy(md_str):
    """
    将Markdown层级结构转换为树形字典
    实现参考知乎专栏《从零构建自己的MCP Server》的AST解析思路
    """
    hierarchy = {}
    current_path = []
    id_counter = 1
    stack = [{'level': 0, 'node': None}]
    
    for line in md_str.strip().split('\n'):
        match = re.match(r'^(#+)\s*(.*)$', line)
        if not match:
            continue
            
        level = len(match.group(1))
        name = match.group(2).strip()
        node_id = f"ID_{uuid.uuid4().hex[:6]}"
        
        # 维护层级栈
        while stack[-1]['level'] >= level:
            stack.pop()
            
        parent = stack[-1]['node']
        node = {
            'id': node_id,
            'name': name,
            'children': []
        }
        
        if parent is None:
            hierarchy[node_id] = node
        else:
            parent['children'].append(node)
            
        stack.append({'level': level, 'node': node})
        
    return hierarchy

def give_hierarchy_positions(hierarchy):
    """
    为每个节点分配位置和大小，实现递归布局算法
    参考知乎《递归布局中的动态宽度计算》和论文《Adaptive Tree Layout with Dynamic Spacing》
    """
    # 布局参数配置（单位：像素）
    # 从配置文件加载布局参数（单位：像素）
    LAYOUT_CONFIG = {
        'canvas_width': config['layout']['canvas_width'],
        'margin': config['layout']['margin'],
        'horizontal_gap': config['layout']['horizontal_gap'],
        'vertical_gap': config['layout']['vertical_gap'],
        'level_height': {
            1: config['level_height']['1'],
            2: config['level_height']['2'],
            3: config['level_height']['3']
        }
    }

    # 转换像素到PPT单位（1英寸=914400 EMU，假设1像素=12700 EMU）
    def px_to_emu(px):
        return int(px * 12700)

    # 递归布局核心算法
    def layout_node(node, parent=None, level=1, start_x=0, start_y=0):
        nonlocal LAYOUT_CONFIG
        
        # 初始化节点布局属性
        node['position'] = {}
        node['size'] = {}
        
        if level == 1:
            # 一级元素布局规则
            node['size']['width'] = LAYOUT_CONFIG['canvas_width'] - 2*LAYOUT_CONFIG['margin']
            node['size']['height'] = LAYOUT_CONFIG['level_height'][1]
            node['position']['x'] = LAYOUT_CONFIG['margin']
            node['position']['y'] = start_y
            
            # 更新全局Y坐标
            max_child_y = start_y
            child_start_x = node['position']['x']
            for idx, child in enumerate(node.get('children', [])):
                child_y = layout_node(child, node, 2, child_start_x, node['position']['y'] + node['size']['height'] + LAYOUT_CONFIG['vertical_gap'])
                max_child_y = max(max_child_y, child_y)
                child_start_x += child['size']['width'] + LAYOUT_CONFIG['horizontal_gap']
            
            node['total_height'] = max_child_y - start_y - LAYOUT_CONFIG['vertical_gap']
            # 为节点增加bounding_box
            node['bounding_box'] = {
                'x': node['position']['x']-2*3,
                'y': node['position']['y']-2*3,
                'width': node['size']['width']+4*3,
                'height': node['total_height']+4*3 
            }
            return max_child_y
            
        elif level == 2:
            # 二级元素动态宽度计算
            siblings_count = len(parent['children'])
            available_width = parent['size']['width'] - (siblings_count-1)*LAYOUT_CONFIG['horizontal_gap']
            node['size']['width'] = available_width / siblings_count
            node['size']['height'] = LAYOUT_CONFIG['level_height'][2]
            
            # 位置计算
            node['position']['x'] = start_x
            node['position']['y'] = start_y
            
            # 处理三级元素
            child_start_x = node['position']['x']
            left =True
            current_child_x = child_start_x
            current_child_y = start_y + node['size']['height'] + LAYOUT_CONFIG['vertical_gap']
            for child in node.get('children', []):
                current_child_y = layout_node(child, node, 3, current_child_x, current_child_y)
                if left:
                    current_child_x += child['size']['width'] + LAYOUT_CONFIG['horizontal_gap']
                    left = False
                else:
                    current_child_y = current_child_y + child['size']['height'] + LAYOUT_CONFIG['vertical_gap']
                    left = True 
            
            node['total_height'] = current_child_y - start_y - LAYOUT_CONFIG['vertical_gap']
            # 为节点增加bounding_box
            node['bounding_box'] = {
                'x': node['position']['x']-2*2,
                'y': node['position']['y']-2*2,
                'width': node['size']['width']+4*2,
                'height': node['total_height']+4*2
            }
            
            return current_child_y
            
        elif level == 3:
            # 三级元素布局 - 改为两两并排
            siblings_count = len(parent['children'])
            available_width = parent['size']['width'] - LAYOUT_CONFIG['horizontal_gap']
            node['size']['width'] = available_width / 2
            node['size']['height'] = LAYOUT_CONFIG['level_height'][3]
            
            node['position']['x'] = start_x
            node['position']['y'] = start_y
            
        return start_y
            
    # 遍历根节点进行布局
    y_offset = LAYOUT_CONFIG['margin']
    for root in [n for n in hierarchy.values() if not any(n in p['children'] for p in hierarchy.values())]:
        y_offset = layout_node(root, level=1, start_y=y_offset) + LAYOUT_CONFIG['vertical_gap']
    
    # 添加单位转换
    def convert_units(node):
        for key in ['x', 'y']:
            node['position'][key] = px_to_emu(node['position'][key])
        for dim in ['width', 'height']:
            node['size'][dim] = px_to_emu(node['size'][dim])
        if 'bounding_box' in node:
            for key in ['x', 'y', 'width', 'height']:
                node['bounding_box'][key] = px_to_emu(node['bounding_box'][key])
        for child in node.get('children', []):
            convert_units(child)
    
    for root in hierarchy.values():
        convert_units(root)
    
    return hierarchy
    
def generate_pptx(hierarchy, style_rules=None):
    """
    生成PPTX架构图，集成知乎《系统架构图绘制规范》样式标准
    参数：

    - hierarchy: 带布局信息的树形结构

    - style_rules: 样式配置字典（可选）
    """
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # 空白版式
    slide = prs.slides.add_slide(slide_layout)
    
    # 加载主题样式
    theme_path = os.path.join(os.path.dirname(__file__), 'theme', f"{config['theme']['template']}.toml")
    with open(theme_path, 'r') as f:
        theme_config = toml.load(f)
    
    # 转换配置为程序可用的格式
    DEFAULT_STYLE = {
        'shape_type': getattr(MSO_SHAPE, theme_config['shape']['type']),
        'level_styles': {
            1: {
                'fill_color': RGBColor(*theme_config['level_styles']['1']['fill_color']),
                'text_color': RGBColor(*theme_config['level_styles']['1']['text_color']),
                'font_size': Pt(theme_config['level_styles']['1']['font_size'])
            },
            2: {
                'fill_color': RGBColor(*theme_config['level_styles']['2']['fill_color']),
                'text_color': RGBColor(*theme_config['level_styles']['2']['text_color']),
                'font_size': Pt(theme_config['level_styles']['2']['font_size'])
            },
            3: {
                'fill_color': RGBColor(*theme_config['level_styles']['3']['fill_color']),
                'text_color': RGBColor(*theme_config['level_styles']['3']['text_color']),
                'font_size': Pt(theme_config['level_styles']['3']['font_size'])
            }
        },
        'bounding_style': {
            'shape_type': getattr(MSO_SHAPE, theme_config['bounding_style']['type']),
            'color': RGBColor(*theme_config['bounding_style']['color']),
            'line_width': Pt(theme_config['bounding_style']['line_width']),
            'dash_style': getattr(MSO_LINE, theme_config['bounding_style']['dash_style'])
        }
    }
    
    # 合并自定义样式
    if style_rules:
        from deepmerge import always_merger
        DEFAULT_STYLE = always_merger.merge(DEFAULT_STYLE, style_rules)
    
    # 形状缓存字典（用于连接线计算）
    node_shapes = {}
    
    def create_shape(node, level):
        """创建形状并设置样式"""
        style = DEFAULT_STYLE['level_styles'][level]
        # 创建基础形状
        shape = slide.shapes.add_shape(
            DEFAULT_STYLE['shape_type'],
            node['position']['x'],
            node['position']['y'],
            node['size']['width'],
            node['size']['height']
        )
        
        # 填充颜色
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = style['fill_color']
        
        # 边框样式
        line = shape.line
        line.color.rgb = RGBColor(0,0,0)
        line.fill.background()
        
        # 文字格式
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = node['name']
        font = run.font
        font.color.rgb = style['text_color']
        font.size = style['font_size']
        font.bold = (level == 1)
        
        return shape

    def draw_bounding(bounding):
        """绘制边界框"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bounding['x'],
            bounding['y'],
            bounding['width'],
            bounding['height']
        )
        style = DEFAULT_STYLE['bounding_style']
        line = shape.line
        line.color.rgb = style['color']
        line.width = style['line_width']
        line.dash_style = style['dash_style']

        fill = shape.fill
        fill.background()
        shape.shadow

        return shape
    
    def draw_node(node, level=1):
        """递归绘制节点及子节点"""
        if 'bounding_box' in node:
            bounding_shape = draw_bounding(node['bounding_box'])
        # 创建当前节点形状
        shape = create_shape(node, level)
        node_shapes[node['id']] = shape

        # 绘制子节点
        for child in node.get('children', []):
            child_shape = draw_node(child, level+1)
        
        return shape

    
    # 遍历根节点开始绘制
    for root in hierarchy.values():
        if not any(root in p['children'] for p in hierarchy.values()):
            draw_node(root)
    
    # 自动调整幻灯片尺寸（根据知乎《PPTX画布动态适配方案》）
    max_x = max(shape.left + shape.width for shape in slide.shapes)
    max_y = max(shape.top + shape.height for shape in slide.shapes)
    prs.slide_width = max_x + Inches(1)
    prs.slide_height = max_y + Inches(1)
    
    return prs

# MCP服务封装（参考知乎《MCP实战-本地MCP Server+Cursor实践》）
mcp = FastMCP("pptx-generator")

@mcp.tool()
def generate_architecture_diagram(md_input: str, style_spec: str = "") -> bytes:
    """
    MCP工具方法：输入markdown和样式规范，输出PPTX字节流
    """
    hierarchy = markdown_to_hierarchy(md_input)
    hierarchy = give_hierarchy_positions(hierarchy)
    prs = generate_pptx(hierarchy, style_spec)
    slide = prs.slides[0]
    for shape in slide.shapes:
        shape.shadow.inherit = False
    # 保存到按照时间戳命名的文件
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"architecture_{timestamp}.pptx"
    save_path = os.path.join(workspace_path, filename)
    prs.save(save_path)
    return f"PPTX文件已保存到：{save_path}"

if __name__ == "__main__":
    # 测试用例（参考用户提供的示例）
    test_md = '''# 管理层
## 策略配置中心
### 资源调度策略库
### 算法参数工作台

## 可视化监控
### 3D拓扑呈现
### 故障传播沙盘

# 应用服务层
## 业务保障系统
### 确定性时延保障
### 业务链自愈

## 网络自治服务
### 智能运维中枢
### 切片编排器

# 智能引擎层
## 算法服务集群
### 图神经网络服务
### 图优化算法库

## 图谱构建引擎
### 动态拓扑建模
### 超图生成器

# 基础设施层
## 计算资源池
### 分布式图计算引擎
### GNN训练平台

## 数据采集系统
### 多源异构数据接入
### 图数据库集群'''
    # 生成PPTX并保存
    print(generate_architecture_diagram(test_md))
